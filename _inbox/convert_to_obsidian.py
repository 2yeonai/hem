from __future__ import annotations

import argparse
import collections
import html
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "옵시디언 변환본"
FULLTEXT = OUTPUT / "_원문"
SUPPORTED_DOCS = {".hwp", ".hwpx", ".docx", ".pptx", ".xlsx", ".pdf"}
SUPPORTED_IMAGES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
SUPPORTED_MEDIA = {".m4a", ".mp3", ".wav", ".mp4", ".webm", ".mov"}
DIRECT_READABLE = SUPPORTED_IMAGES | SUPPORTED_MEDIA | {".pdf", ".md"}


def vault_path(path: Path) -> str:
    return path.resolve().relative_to(ROOT).as_posix()


def clean_text(text: str) -> str:
    text = text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n[ \t]+", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def summarize(text: str, limit: int = 8) -> tuple[list[str], list[str]]:
    compact = clean_text(text)
    words = re.findall(r"[가-힣A-Za-z0-9]{2,}", compact.lower())
    stop = {
        "그리고", "그러나", "대한", "위한", "있는", "하는", "합니다", "있습니다", "됩니다",
        "the", "and", "for", "with", "from", "this", "that", "are", "was", "were",
    }
    freq = collections.Counter(w for w in words if w not in stop and not w.isdigit())
    keywords = [w for w, _ in freq.most_common(15)]
    sentences = [s.strip(" -•\t") for s in re.split(r"(?<=[.!?다요함됨음])\s+|\n+", compact) if len(s.strip()) >= 20]
    scored = []
    for idx, sentence in enumerate(sentences):
        tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", sentence.lower())
        score = sum(freq.get(t, 0) for t in set(tokens)) / max(len(tokens), 1)
        if len(sentence) > 500:
            score *= 0.6
        scored.append((score, idx, sentence[:500]))
    chosen = sorted(sorted(scored, reverse=True)[:limit], key=lambda x: x[1])
    key_sentences = [item[2] for item in chosen]
    if not key_sentences and compact:
        key_sentences = [compact[:1000]]
    return key_sentences, keywords


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    chunks: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)
    for table in doc.tables:
        for row in table.rows:
            values = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            if any(values):
                chunks.append(" | ".join(values))
    return clean_text("\n\n".join(chunks))


def extract_hwpx(path: Path) -> str:
    chunks: list[str] = []
    with zipfile.ZipFile(path) as archive:
        names = sorted(
            (n for n in archive.namelist() if re.fullmatch(r"Contents/section\d+\.xml", n)),
            key=lambda n: int(re.search(r"\d+", n).group()),
        )
        for name in names:
            root = ET.fromstring(archive.read(name))
            for elem in root.iter():
                if elem.tag.rsplit("}", 1)[-1] == "p":
                    parts = [node.text or "" for node in elem.iter() if node.tag.rsplit("}", 1)[-1] == "t"]
                    line = clean_text("".join(parts))
                    if line:
                        chunks.append(line)
    return clean_text("\n\n".join(chunks))


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    chunks: list[str] = []
    for number, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = clean_text(shape.text)
                if value:
                    lines.append(value)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
                    if any(values):
                        lines.append(" | ".join(values))
        chunks.append(f"## 슬라이드 {number}\n\n" + "\n\n".join(lines))
    return clean_text("\n\n".join(chunks))


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    book = load_workbook(path, data_only=True, read_only=True)
    chunks: list[str] = []
    for sheet in book.worksheets:
        chunks.append(f"## 시트: {sheet.title}")
        for row in sheet.iter_rows():
            values = [str(cell.value).strip() if cell.value is not None else "" for cell in row]
            if any(values):
                chunks.append(" | ".join(values))
    book.close()
    return clean_text("\n\n".join(chunks))


def extract_hwp(path: Path) -> str:
    import win32com.client

    app = win32com.client.gencache.EnsureDispatch("HwpFrame.HwpObject")
    try:
        try:
            app.XHwpWindows.Item(0).Visible = False
        except Exception:
            pass
        if not app.Open(str(path), "HWP", "forceopen:true"):
            raise RuntimeError("한컴오피스에서 문서를 열지 못했습니다")
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as handle:
            temp_path = Path(handle.name)
        try:
            if not app.SaveAs(str(temp_path), "TEXT", ""):
                raise RuntimeError("한컴오피스 텍스트 저장에 실패했습니다")
            data = temp_path.read_bytes()
            for encoding in ("utf-16", "utf-8-sig", "cp949"):
                try:
                    return clean_text(data.decode(encoding))
                except UnicodeDecodeError:
                    continue
            return clean_text(data.decode("utf-8", errors="replace"))
        finally:
            temp_path.unlink(missing_ok=True)
    finally:
        try:
            app.Clear(1)
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass


def extract_pdf(path: Path) -> str:
    import fitz

    document = fitz.open(path)
    chunks: list[str] = []
    for number, page in enumerate(document, 1):
        text = clean_text(page.get_text("text"))
        chunks.append(f"## 페이지 {number}\n\n{text or '[텍스트가 없는 이미지형 페이지 — OCR 확인 필요]'}")
    document.close()
    return clean_text("\n\n".join(chunks))


def get_ocr_engine():
    try:
        from rapidocr_onnxruntime import RapidOCR

        return RapidOCR()
    except Exception:
        return None


def extract_image(path: Path, ocr_engine) -> str:
    if ocr_engine is None:
        return "[이미지 문자 인식 도구가 준비되지 않았습니다]"
    result, _ = ocr_engine(str(path))
    if not result:
        return "[이미지에서 인식된 문자가 없습니다]"
    return clean_text("\n".join(str(item[1]) for item in result))


def transcribe_media(path: Path, model: dict[str, Path]) -> str:
    import av
    import imageio_ffmpeg

    safe_name = re.sub(r"[^0-9A-Za-z가-힣_-]+", "_", path.stem)[:80]
    temp = ROOT / "work" / "media-temp" / safe_name
    checkpoints = ROOT / "work" / "media-checkpoints" / safe_name
    temp.mkdir(parents=True, exist_ok=True)
    checkpoints.mkdir(parents=True, exist_ok=True)
    try:
        ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
        container = av.open(str(path))
        duration = float(container.duration / av.time_base) if container.duration else 0
        container.close()
        if duration <= 0:
            raise RuntimeError("미디어 재생시간을 확인하지 못했습니다")
        threads = max(2, os.cpu_count() or 4)
        chunk_seconds = 30 * 60
        chunks: list[str] = []
        for chunk_no, offset in enumerate(range(0, int(duration) + 1, chunk_seconds), 1):
            if offset >= duration:
                break
            length = min(chunk_seconds, duration - offset)
            checkpoint = checkpoints / f"chunk-{chunk_no:02d}.txt"
            if checkpoint.exists():
                raw = clean_text(checkpoint.read_text(encoding="utf-8", errors="replace"))
                stamp = f"{offset // 3600:02d}:{(offset % 3600) // 60:02d}:{offset % 60:02d}"
                chunks.append(f"## 구간 시작 {stamp}\n\n{raw}")
                continue
            wav_path = temp / f"audio-{chunk_no:02d}.wav"
            output_base = temp / f"transcript-{chunk_no:02d}"
            subprocess.run(
                [
                    ffmpeg, "-y", "-ss", str(offset), "-t", str(length), "-i", str(path),
                    "-vn", "-ar", "16000", "-ac", "1", "-c:a", "pcm_s16le", str(wav_path),
                ],
                check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            )
            command = [
                str(model["cli"]), "-m", str(model["model"]), "-f", str(wav_path),
                "-l", "ko", "-t", str(threads), "-otxt", "-of", str(output_base),
                "-bs", "1", "-bo", "1", "-nf", "-ng", "--no-prints",
            ]
            subprocess.run(command, check=True)
            transcript_path = Path(str(output_base) + ".txt")
            raw = clean_text(transcript_path.read_text(encoding="utf-8", errors="replace"))
            checkpoint.write_text(raw, encoding="utf-8")
            stamp = f"{offset // 3600:02d}:{(offset % 3600) // 60:02d}:{offset % 60:02d}"
            chunks.append(f"## 구간 시작 {stamp}\n\n{raw}")
            wav_path.unlink(missing_ok=True)
            transcript_path.unlink(missing_ok=True)
        return clean_text("\n\n".join(chunks))
    finally:
        shutil.rmtree(temp, ignore_errors=True)


def load_whisper(model_size: str) -> dict[str, Path]:
    model_files = {
        "base": "ggml-base.bin",
        "small": "ggml-small.bin",
        "small-q5_1": "ggml-small-q5_1.bin",
    }
    if model_size not in model_files:
        raise ValueError("현재 준비된 무료 로컬 모델은 base, small 또는 small-q5_1입니다")
    base = ROOT / "work" / "whispercpp"
    cli = base / "bin" / "Release" / "whisper-cli.exe"
    model = base / model_files[model_size]
    if not cli.exists() or not model.exists():
        raise FileNotFoundError("Windows용 Whisper 실행판 또는 모델이 없습니다")
    return {"cli": cli, "model": model}


def extract_document(path: Path) -> str:
    dispatch = {
        ".hwp": extract_hwp,
        ".hwpx": extract_hwpx,
        ".docx": extract_docx,
        ".pptx": extract_pptx,
        ".xlsx": extract_xlsx,
        ".pdf": extract_pdf,
    }
    return dispatch[path.suffix.lower()](path)


def write_note(source: Path, full_text: str, method: str, warnings: list[str] | None = None) -> tuple[Path, Path]:
    relative = source.relative_to(ROOT)
    note_path = OUTPUT / relative.parent / f"{source.stem}.md"
    full_path = FULLTEXT / relative.parent / f"{source.stem}__원문.md"
    note_path.parent.mkdir(parents=True, exist_ok=True)
    full_path.parent.mkdir(parents=True, exist_ok=True)
    key_sentences, keywords = summarize(full_text)
    source_link = vault_path(source)
    full_link = vault_path(full_path)
    warning_lines = "\n".join(f"> - {item}" for item in (warnings or []))
    note = f'''---
title: "{source.stem.replace(chr(34), chr(39))}"
type: converted-document
status: converted
source_format: {source.suffix.lower().lstrip('.')}
source: "{source_link}"
converted: 2026-07-13
tags: [문서변환, 인박스, 옵시디언]
---

# {source.stem}

> [!info] 변환 정보
> 변환 방식: {method}  
> [[{source_link}|원본 파일]] · [[{full_link}|추출 원문 전체]]
{warning_lines}

## AI 빠른 읽기

'''
    note += "\n".join(f"- {line}" for line in key_sentences)
    note += "\n\n## 핵심어\n\n" + (", ".join(f"#{word}" for word in keywords) or "핵심어 추출 없음") + "\n"
    full_note = f'''---
title: "{source.stem.replace(chr(34), chr(39))} - 추출 원문"
type: extracted-fulltext
source: "{source_link}"
converted: 2026-07-13
tags: [문서변환, 원문, 인박스]
---

# {source.stem} - 추출 원문

[[{vault_path(note_path)}|AI 빠른 읽기 문서로 돌아가기]] · [[{source_link}|원본 파일]]

{full_text}
'''
    note_path.write_text(note, encoding="utf-8")
    full_path.write_text(full_note, encoding="utf-8")
    return note_path, full_path


def source_files() -> list[Path]:
    excluded = {OUTPUT.resolve(), (ROOT / "work").resolve(), (ROOT / ".git").resolve()}
    files: list[Path] = []
    for path in ROOT.rglob("*"):
        if not path.is_file():
            continue
        resolved = path.resolve()
        if any(root == resolved or root in resolved.parents for root in excluded):
            continue
        if path.name in {Path(__file__).name, "convert_to_obsidian.ps1"}:
            continue
        files.append(path)
    return sorted(files, key=lambda p: vault_path(p).lower())


def write_index(results: list[dict], all_files: list[Path]) -> None:
    converted = {item["source"].resolve(): item for item in results}
    for source in all_files:
        if source.resolve() in converted:
            continue
        relative = source.relative_to(ROOT)
        existing_note = OUTPUT / relative.parent / f"{source.stem}.md"
        existing_full = FULLTEXT / relative.parent / f"{source.stem}__원문.md"
        if existing_note.exists():
            converted[source.resolve()] = {
                "source": source, "status": "완료", "note": existing_note,
                "full": existing_full if existing_full.exists() else None, "error": "",
            }
    lines = [
        "---", 'title: "인박스 자료 인덱스"', "type: index", "status: active",
        "updated: 2026-07-13", "tags: [문서변환, 인박스, 옵시디언]", "---", "",
        "# 인박스 자료 인덱스", "",
        "먼저 변환 문서의 `AI 빠른 읽기`를 보고, 세부 확인이 필요할 때만 `추출 원문 전체`를 여세요.", "",
    ]
    groups: dict[str, list[Path]] = collections.defaultdict(list)
    for source in all_files:
        parent = str(source.relative_to(ROOT).parent)
        groups["인박스 루트" if parent == "." else parent].append(source)
    for group in sorted(groups):
        lines.extend([f"## {group}", ""])
        for source in groups[group]:
            item = converted.get(source.resolve())
            src = vault_path(source)
            if item and item["status"] == "완료":
                lines.append(f"- [[{vault_path(item['note'])}|{source.name}]] — 변환 완료 · [[{src}|원본]]")
            elif item:
                lines.append(f"- **{source.name}** — 변환 실패: {item['error']} · [[{src}|원본]]")
            elif source.suffix.lower() in DIRECT_READABLE:
                lines.append(f"- [[{src}|{source.name}]] — 옵시디언에서 원본 열기")
            else:
                lines.append(f"- [[{src}|{source.name}]] — 원본 연결")
        lines.append("")
    counts = collections.Counter(item["status"] for item in results)
    lines.extend(["## 처리 결과", ""] + [f"- {key}: {value}개" for key, value in sorted(counts.items())])
    OUTPUT.mkdir(parents=True, exist_ok=True)
    (OUTPUT / "00 인박스 자료 인덱스.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--phase", choices=("docs", "images", "media", "all"), default="docs")
    parser.add_argument("--whisper-model", default="small")
    parser.add_argument("--match", help="파일 경로에 이 문자열이 포함된 대상만 처리")
    parser.add_argument("--skip-existing", action="store_true")
    args = parser.parse_args()
    all_files = source_files()
    targets: list[Path] = []
    if args.phase in {"docs", "all"}:
        targets.extend(p for p in all_files if p.suffix.lower() in SUPPORTED_DOCS)
    if args.phase in {"images", "all"}:
        targets.extend(p for p in all_files if p.suffix.lower() in SUPPORTED_IMAGES)
    if args.phase in {"media", "all"}:
        targets.extend(p for p in all_files if p.suffix.lower() in SUPPORTED_MEDIA)
    if args.match:
        targets = [p for p in targets if args.match.lower() in vault_path(p).lower()]
    if args.skip_existing:
        targets = [
            p for p in targets
            if not (OUTPUT / p.relative_to(ROOT).parent / f"{p.stem}.md").exists()
        ]

    ocr_engine = get_ocr_engine() if args.phase in {"images", "all"} else None
    whisper_model = load_whisper(args.whisper_model) if any(p.suffix.lower() in SUPPORTED_MEDIA for p in targets) else None
    results: list[dict] = []
    for number, source in enumerate(targets, 1):
        print(f"[{number}/{len(targets)}] {vault_path(source)}", flush=True)
        try:
            suffix = source.suffix.lower()
            if suffix in SUPPORTED_DOCS:
                text = extract_document(source)
                method = "로컬 문서 본문 추출"
                warnings = ["복잡한 표·도형·서식은 단순화될 수 있습니다."]
            elif suffix in SUPPORTED_IMAGES:
                text = extract_image(source, ocr_engine)
                method = "로컬 이미지 문자 인식(OCR)"
                warnings = ["문자 중심 결과이며 사진의 시각적 의미 설명은 포함하지 않습니다."]
            else:
                text = transcribe_media(source, whisper_model)
                method = "로컬 음성 인식 전사"
                warnings = ["고유명사·숫자는 오인식될 수 있어 원본 확인이 필요합니다."]
            note, full = write_note(source, text, method, warnings)
            results.append({"source": source, "status": "완료", "note": note, "full": full, "error": ""})
        except Exception as exc:
            results.append({"source": source, "status": "실패", "note": None, "full": None, "error": str(exc)})
            print(f"  실패: {exc}", file=sys.stderr, flush=True)
    write_index(results, all_files)
    counts = collections.Counter(item["status"] for item in results)
    for key, value in sorted(counts.items()):
        print(f"{key}: {value}")
    return 0 if not counts.get("실패") else 1


if __name__ == "__main__":
    raise SystemExit(main())
